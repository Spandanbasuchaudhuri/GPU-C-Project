"""
ingest.py  –  extract text / OCR / ASR
"""
import io, os, torch, pdfplumber, pytesseract, config
from PIL import Image

# optional deps
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
try:
    from pdf2image import convert_from_path
except ImportError:
    convert_from_path = None

# ── OCR ──────────────────────────────────────────────────────────────────
def _ocr(img: Image.Image) -> str:
    return pytesseract.image_to_string(img, lang=config.TESSERACT_LANG).strip()

# ── PDF ──────────────────────────────────────────────────────────────────
def extract_text_pdf(path: str) -> str:
    out = []
    with pdfplumber.open(path) as pdf:
        for n, page in enumerate(pdf.pages, 1):
            txt = (page.extract_text() or "").strip()
            if txt:
                out.append(txt); continue
            try:
                img = page.to_image(resolution=300).original
                if (o := _ocr(img)):
                    out.append(f"[OCR-p{n}] {o}")
            except Exception:
                pass

    if not out and convert_from_path:  # encrypted PDFs
        for img in convert_from_path(path, dpi=300):
            if (o := _ocr(img)):
                out.append(f"[OCR] {o}")
    return "\n".join(out)

# ── PPTX ─────────────────────────────────────────────────────────────────
def extract_text_pptx(path: str) -> str:
    if Presentation is None:
        raise ImportError("pip install python-pptx")
    prs, out = Presentation(path), []
    for sl in prs.slides:
        for sh in sl.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                out.append(sh.text_frame.text.strip())
            elif getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        out.append(" | ".join(cells))
    return "\n".join(out)

# ── Images ───────────────────────────────────────────────────────────────
def ocr_image(path: str) -> str:
    return _ocr(Image.open(path).convert("RGB"))

# ── Audio ────────────────────────────────────────────────────────────────
_whisper = None
def transcribe_audio(path: str) -> str:
    global _whisper
    if _whisper is None:
        from faster_whisper import WhisperModel       # pip install faster-whisper
        dev  = "cuda" if torch.cuda.is_available() else "cpu"
        comp = "float16" if dev == "cuda" else "int8"
        _whisper = WhisperModel(config.WHISPER_MODEL, device=dev, compute_type=comp)
    segs, _ = _whisper.transcribe(path, beam_size=5, vad_filter=True, language="en")
    return " ".join(s.text for s in segs)